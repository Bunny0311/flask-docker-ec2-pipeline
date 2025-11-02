from flask import Flask, render_template, request, redirect, url_for, send_file
import os
from werkzeug.utils import secure_filename
import pdfplumber
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import io

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
CONVERTED_FOLDER = "converted"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER

# ---------- Helper Functions ----------

def pdf_to_word(pdf_path, output_path):
    with pdfplumber.open(pdf_path) as pdf:
        doc = Document()
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                doc.add_paragraph(text)
        doc.save(output_path)

def word_to_pdf(word_path, output_path):
    doc = Document(word_path)
    pdf = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    y = height - 50
    for para in doc.paragraphs:
        pdf.drawString(50, y, para.text)
        y -= 15
        if y < 50:
            pdf.showPage()
            y = height - 50
    pdf.save()

def pdf_to_ppt(pdf_path, output_path):
    prs = Presentation()
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Converted from PDF"
            slide.placeholders[1].text = text if text else ""
    prs.save(output_path)

def ppt_to_pdf(ppt_path, output_path):
    prs = Presentation(ppt_path)
    pdf = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    for slide in prs.slides:
        y = height - 50
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.drawString(50, y, shape.text)
                y -= 15
        pdf.showPage()
    pdf.save()

def excel_to_csv(excel_path, output_path):
    df = pd.read_excel(excel_path)
    df.to_csv(output_path, index=False)

def csv_to_excel(csv_path, output_path):
    df = pd.read_csv(csv_path)
    df.to_excel(output_path, index=False)

# ---------- Routes ----------

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return "No file uploaded"

    file = request.files['file']
    convert_type = request.form.get('convert_type')

    if file.filename == '':
        return "No file selected"

    filename = secure_filename(file.filename)
    upload_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(upload_path)

    base_name, ext = os.path.splitext(filename)
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], base_name)

    # Conversion logic
    try:
        if convert_type == 'pdf_to_word':
            output_file = output_path + ".docx"
            pdf_to_word(upload_path, output_file)

        elif convert_type == 'word_to_pdf':
            output_file = output_path + ".pdf"
            word_to_pdf(upload_path, output_file)

        elif convert_type == 'pdf_to_ppt':
            output_file = output_path + ".pptx"
            pdf_to_ppt(upload_path, output_file)

        elif convert_type == 'ppt_to_pdf':
            output_file = output_path + ".pdf"
            ppt_to_pdf(upload_path, output_file)

        elif convert_type == 'excel_to_csv':
            output_file = output_path + ".csv"
            excel_to_csv(upload_path, output_file)

        elif convert_type == 'csv_to_excel':
            output_file = output_path + ".xlsx"
            csv_to_excel(upload_path, output_file)

        else:
            return "Unsupported conversion type"

    except Exception as e:
        return f"Error: {e}"

    return render_template('result.html', output_file=os.path.basename(output_file))

@app.route('/download/<filename>')
def download_file(filename):
    return send_file(os.path.join(app.config['CONVERTED_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
